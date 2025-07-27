"""
Document processing service using LangGraph for parsing and chunking documents.
"""
import os
import csv
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
from io import StringIO

import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document as LangChainDocument

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Service for processing various document types."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: Size of text chunks for splitting
            chunk_overlap: Overlap between chunks
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    def process_file(self, file_path: str, filename: str) -> List[LangChainDocument]:
        """
        Process a single file and return chunked documents.
        
        Args:
            file_path: Path to the file
            filename: Original filename
            
        Returns:
            List of LangChain Document objects
        """
        try:
            # Extract text based on file extension
            file_ext = Path(filename).suffix.lower()
            
            if file_ext == '.pdf':
                text = self._extract_pdf_text(file_path)
            elif file_ext == '.docx':
                text = self._extract_docx_text(file_path)
            elif file_ext == '.txt':
                text = self._extract_txt_text(file_path)
            elif file_ext == '.pptx':
                text = self._extract_pptx_text(file_path)
            elif file_ext == '.csv':
                text = self._extract_csv_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_ext}")
            
            if not text.strip():
                logger.warning(f"No text extracted from {filename}")
                return []
            
            # Create a single document with metadata
            document = LangChainDocument(
                page_content=text,
                metadata={
                    "source": filename,
                    "file_path": file_path,
                    "file_type": file_ext,
                    "chunk_index": 0
                }
            )
            
            # Split the document into chunks
            chunks = self.text_splitter.split_documents([document])
            
            # Update metadata for each chunk
            for i, chunk in enumerate(chunks):
                chunk.metadata.update({
                    "chunk_index": i,
                    "total_chunks": len(chunks)
                })
            
            logger.info(f"Processed {filename}: {len(chunks)} chunks created")
            return chunks
            
        except Exception as e:
            logger.error(f"Error processing file {filename}: {str(e)}")
            raise
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file."""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                    except Exception as e:
                        logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                        continue
        except Exception as e:
            logger.error(f"Error reading PDF file: {e}")
            raise
        
        return text.strip()
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        try:
            doc = Document(file_path)
            text_parts = []
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"Error reading DOCX file: {e}")
            raise
    
    def _extract_txt_text(self, file_path: str) -> str:
        """Extract text from TXT file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                logger.error(f"Error reading TXT file with latin-1 encoding: {e}")
                raise
        except Exception as e:
            logger.error(f"Error reading TXT file: {e}")
            raise
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        try:
            prs = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip() != f"--- Slide {slide_num + 1} ---":
                    text_parts.append(slide_text)
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"Error reading PPTX file: {e}")
            raise
    
    def _extract_csv_text(self, file_path: str) -> str:
        """Extract text from CSV file."""
        try:
            # Try to read with pandas first
            df = pd.read_csv(file_path)
            
            # Convert DataFrame to a readable text format
            text_parts = []
            text_parts.append(f"CSV File with {len(df)} rows and {len(df.columns)} columns")
            text_parts.append(f"Columns: {', '.join(df.columns.tolist())}")
            text_parts.append("\nData Summary:")
            text_parts.append(df.describe(include='all').to_string())
            
            # Add sample rows
            text_parts.append(f"\nFirst 10 rows:")
            text_parts.append(df.head(10).to_string())
            
            return "\n".join(text_parts)
            
        except Exception as e:
            logger.error(f"Error reading CSV file with pandas: {e}")
            # Fallback to basic CSV reading
            try:
                with open(file_path, 'r', encoding='utf-8') as file:
                    csv_reader = csv.reader(file)
                    rows = list(csv_reader)
                    
                text_parts = []
                if rows:
                    text_parts.append(f"CSV File with {len(rows)} rows")
                    if len(rows) > 0:
                        text_parts.append(f"Headers: {', '.join(rows[0])}")
                    
                    # Add first few rows
                    for i, row in enumerate(rows[:11]):  # Header + 10 data rows
                        text_parts.append(" | ".join(row))
                
                return "\n".join(text_parts)
                
            except Exception as e2:
                logger.error(f"Error reading CSV file with basic reader: {e2}")
                raise

